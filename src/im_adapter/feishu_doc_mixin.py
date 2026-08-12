"""飞书文档能力混入（Mixin）—— 文档搜索 / 读取 / 列表 / 全量消息拉取。

混入 ``FeishuCliAdapter``，通过 MRO 访问 ``BaseIMAdapter`` 提供的
``run()`` / ``_build_command()`` / ``_run_download()`` 以及自身类本体提供的
``_normalize_message()`` / ``_infer_single_chat()`` / ``chat_conversation_info()``
等核心依赖。
"""
from __future__ import annotations
from .im_mixins_base import IMAdapterBase

import json
import logging
import os
import shutil
import subprocess
import tempfile
from typing import Any

from .errors import (
    IMAdapterError,
    IMAdapterPermissionError,
    IMAdapterResourceNotFoundError,
    IMAdapterRateLimitError,
    IMAdapterUnsupportedTypeError,
    IMAdapterShutdownError,
)

logger = logging.getLogger(__name__)


class FeishuDocMixin(IMAdapterBase):
    """飞书文档能力混入 —— 提供文档搜索 / 读取 / 列表 / 全量消息拉取等能力。

    不直接继承任何基类；所有依赖（``self.run()`` / ``self._build_command()`` /
    ``self._no_browser_env`` / ``self.cli_path`` 等）均由最终类
    ``FeishuCliAdapter`` 的 MRO 在运行时提供。
    """

    def chat_message_list_all(self, start: str, end: str,
                              limit: int = 50,
                              max_pages: int | None = None,
                              extra_chat_ids: list[str] | None = None,
                              chat_ids: list[str] | None = None,
                              chat_meta: dict[str, dict] | None = None) -> dict:
        """按时间范围拉取所有消息（单聊 + 群聊），自动分页聚合。

        两种模式：
        - 全量模式（chat_ids 为空，默认）：列出全部会话（p2p+group，分页封顶
          ``MAX_CHATS``），对每个会话翻页拉消息。用于历史同步 / 新会话发现。
        - 白名单模式（chat_ids 非空）：跳过 ``+chat-list`` 全量翻页，仅对传入的
          已知相关会话 + 单次最近活跃探嗅（不翻页）拉消息。用于每轮轮询，
          把 CLI 调用量从「翻页 200 会话 + 逐会话」降到「几十个已知会话」。

        Args:
            start: 查询起始时间（ISO 8601）。
            end: 查询结束时间（ISO 8601）。
            limit: 每页消息数上限。
            extra_chat_ids: 额外强制纳入的 chat_id（兼容旧调用，如外部好友）。
            chat_ids: 白名单模式会话集合；非空即进入白名单模式。
            chat_meta: ``{chat_id: {"title", "chat_mode", "singleChat"}}``，
                为白名单会话提供元数据（避免额外 +chat-list 查询）。

        返回格式（poller 兼容）：
        ``{"conversationMessagesList": [{"openConversationId": ..., "title": ..., "messages": [...], "singleChat": ...}, ...]}``
        """
        MAX_CHATS = 200
        by_chat: dict[str, list[dict]] = {}
        chat_info: dict[str, dict] = {}
        chat_meta = chat_meta or {}

        # 0) 收集需要强制纳入的额外 chat_id（去重）
        extra_ids: set[str] = {cid for cid in (extra_chat_ids or []) if cid}

        if chat_ids:
            # ======== 白名单模式：跳过 +chat-list 全量翻页 ========
            whitelist = list(dict.fromkeys([*chat_ids, *extra_ids]))
            chats: list[dict] = []
            for cid in whitelist:
                if not cid:
                    continue
                meta = chat_meta.get(cid) or {}
                info: dict[str, Any] = {
                    "chat_id": cid,
                    "id": cid,
                    "name": meta.get("title") or meta.get("name") or "",
                    "chat_mode": meta.get("chat_mode")
                    or ("p2p" if meta.get("singleChat") else ""),
                }
                chat_info[cid] = info
                chats.append(info)
            # 单次最近活跃会话探嗅（page-size 封顶、不翻页），
            # 捕获白名单之外本回合新活跃的会话（刚被拉进的群 / 新单聊），避免漏消息。
            try:
                peek = self.run([
                    "im", "+chat-list", "--types", "p2p,group",
                    "--page-size", str(min(30, 50)),
                ])
                for c in self._items(peek):
                    cid = c.get("chat_id") or c.get("id") or ""
                    if cid and cid not in chat_info:
                        chat_info[cid] = c
                        chats.append(c)
            except Exception:  # noqa: BLE001
                logger.warning("[resilience] silent exception in chat_message_list_all", exc_info=True)
            logger.debug(
                "feishu chat_message_list_all: 白名单模式，待拉 %d 个会话"
                "（已跳过 +chat-list 全量翻页）", len(chats))
        else:
            # ======== 全量模式（兼容历史同步 / 新会话发现）========
            # 1) 列出全部会话（分页），同时缓存会话信息
            chats = []
            seen_chat_ids: set[str] = set()
            page_token = ""
            while len(chats) < MAX_CHATS:
                cargs = ["im", "+chat-list", "--types", "p2p,group",
                         "--page-size", "100"]
                if page_token:
                    cargs += ["--page-token", page_token]
                try:
                    cresp = self.run(cargs)
                except Exception:  # noqa: BLE001
                    logger.warning("[resilience] silent exception in chat_message_list_all", exc_info=True)
                    break
                page_items = self._items(cresp)
                for c in page_items:
                    cid = c.get("chat_id") or c.get("id") or ""
                    if cid:
                        chat_info[cid] = c
                        seen_chat_ids.add(cid)
                chats.extend(page_items)
                page_token = (self._payload(cresp) or {}).get("page_token") or ""
                if not page_token or not page_items:
                    break

            # 1.5) 将未出现在 chat-list 中的 extra_chat_ids 注入聊天列表，
            #       确保外部好友等不常联系的会话也被轮询。
            #       通过 chat_conversation_info 补全会话元数据（名称、类型等）。
            #       必须带 chat_mode，否则 _infer_single_chat 会误判为 group。
            for extra_cid in extra_ids:
                if extra_cid in seen_chat_ids:
                    continue
                if len(chats) >= MAX_CHATS:
                    break
                try:
                    info = self.chat_conversation_info(extra_cid)
                except Exception:
                    logger.warning("[resilience] silent exception in chat_message_list_all", exc_info=True)
                    info = {}
                fake_chat: dict[str, Any] = {
                    "chat_id": extra_cid,
                    "id": extra_cid,
                    "name": info.get("name") or info.get("title") or "",
                    "chat_mode": info.get("chat_mode") or "",
                }
                chat_info[extra_cid] = fake_chat
                chats.append(fake_chat)
                seen_chat_ids.add(extra_cid)
                logger.debug(
                    "feishu chat_message_list_all: 强制注入外部好友会话 chat_id=%s chat_mode=%s",
                    extra_cid[:30], fake_chat.get("chat_mode", "?"),
                )

        # 2) 逐会话拉消息（分页）—— 两种模式共享
        blocked_chats: list[dict] = []  # 收集遍历中命中的永久权限错误会话，交由轮询器拉黑
        for chat in chats[:MAX_CHATS]:
            chat_id = chat.get("chat_id") or chat.get("id") or ""
            if not chat_id:
                continue
            msgs: list[dict] = []
            token = ""
            max_pages = 50
            pages = 0
            while True:
                pages += 1
                margs = ["im", "+chat-messages-list", "--chat-id", chat_id,
                         "--order", "asc", "--page-size", str(min(limit, 50)),
                         "--start", start, "--end", end]
                if token:
                    margs += ["--page-token", token]
                try:
                    mresp = self.run(margs)
                except Exception as e:  # noqa: BLE001
                    # 跨租户/跨 app/已退群等永久权限错误属正常业务边界：收集待拉黑，
                    # 交由轮询器写入当前账号黑名单，后续轮询直接跳过、不再遍历消息；
                    # 其余瞬时错误保留 WARNING 便于排查。
                    if isinstance(e, IMAdapterPermissionError):
                        meta = chat_info.get(chat_id, {})
                        title = chat.get("name") or chat.get("title") or meta.get("title", "")
                        chat_mode = chat.get("chat_mode") or meta.get("chat_mode") or ""
                        blocked_chats.append({
                            "chat_id": chat_id,
                            "title": title,
                            "chat_type": "single" if chat_mode == "p2p" else ("group" if chat_mode == "group" else ""),
                            "error": str(e),
                        })
                        logger.debug(
                            "feishu chat_message_list_all: 收集不可达会话待拉黑 chat_id=%s，原因: %s",
                            chat_id, str(e)
                        )
                    else:
                        logger.warning(
                            "feishu chat_message_list_all: 跳过会话 chat_id=%s，原因: %s",
                            chat_id, str(e)
                        )
                    break  # 该会话本轮拉取失败，跳过并进入下一个 chat（run() 自身已负责重试）
                page_items = self._items(mresp)
                msgs.extend([self._normalize_message(m) for m in page_items])
                token = (self._payload(mresp) or {}).get("page_token") or ""
                if not token or not page_items:
                    break
                if pages >= max_pages:
                    logger.warning("分页超过 50 页上限，强制终止")
                    break
            if msgs:
                by_chat[chat_id] = msgs

        # 3) 组装为 poller 兼容格式
        conv_list: list[dict] = []
        for chat_id, msgs in by_chat.items():
            info = chat_info.get(chat_id, {})
            conv_list.append({
                "openConversationId": chat_id,
                "title": info.get("name") or info.get("title") or "",
                "messages": msgs,
                "singleChat": self._infer_single_chat(info),
            })

        return {
            "conversationMessagesList": conv_list,
            "by_chat": by_chat,
            "chat_count": len(by_chat),
            "blocked_chats": blocked_chats,
        }


    # ------------------------------------------------------------------
    # 文档能力
    # ------------------------------------------------------------------

    def _unwrap_wiki_file_token(self, meta: dict) -> str | None:
        """从 wiki 文件节点的 ``result_meta.icon_info`` 取出底层 Drive 文件 token。

        lark-cli ``docs +search`` 对知识库里的文件节点，顶层 ``token`` 是 wiki 节点
        token（直接 ``drive +download`` 会 404），真正的文件 token 在 ``icon_info``
        （JSON 字符串）的 ``token`` 字段里。
        """
        icon_info_raw = meta.get("icon_info")
        if not isinstance(icon_info_raw, str) or not icon_info_raw:
            return None
        try:
            icon_info = json.loads(icon_info_raw)
        except (ValueError, TypeError) as _exc:
            logger.debug(f"_unwrap_wiki_file_token: swallowed exception: {_exc}")
            return None
        tok = icon_info.get("token")
        return tok if isinstance(tok, str) and tok else None

    def doc_search(self, query: str, page_size: int = 10) -> list[dict]:
        """搜索飞书文档 / 知识库内容。

        通过 ``lark-cli docs +search`` 命令（Search v2: doc_wiki/search）搜索
        当前用户有权限访问的文档。
        返回 list[dict]，每项含 doc_token / title / url / snippet 等字段。
        认证失败时返回含 error 标记的 list 而非空 list，以便调用方做 UI 提示。
        """
        args = ["docs", "+search", "--query", query,
                "--page-size", str(page_size)]
        try:
            result = self.run(args, operation="doc_search",
                            force_no_dry_run=True)
            logger.debug("飞书 doc_search 原始返回: %s",
                        json.dumps(result, ensure_ascii=False)[:500])
            if isinstance(result, dict) and result.get("ok"):
                items = result.get("data", [])
                # lark-cli docs +search 实际返回 {"data": {"results": [...], ...}}
                if isinstance(items, dict):
                    raw_list = items.get("results") or items.get("items") or []
                else:
                    raw_list = items if isinstance(items, list) else []
                # 将 lark-cli 原始字段映射为统一格式
                docs = []
                for r in raw_list:
                    if not isinstance(r, dict):
                        continue
                    meta = r.get("result_meta", {})
                    # title: 优先去高亮标题，fallback 普通字段
                    title = r.get("title_highlighted", "") or meta.get("title", "")
                    # snippet
                    snippet = r.get("summary_highlighted", "")
                    token = meta.get("token", "") or r.get("token", "")
                    url = meta.get("url", "")
                    entity_type = r.get("entity_type", "DOC")
                    # 知识库（wiki）里的「文件」节点：搜索返回的顶层 token 是 wiki 节点
                    # token，真正的 Drive 文件 token 藏在 result_meta.icon_info 内。直接用
                    # wiki 节点 token 去 drive +download 会 404，必须换成底层文件 token 才能下载。
                    if entity_type == "WIKI":
                        doc_types = (meta.get("doc_types") or "").upper()
                        file_type = meta.get("file_type") or ""
                        if doc_types == "FILE" or file_type:
                            real_token = self._unwrap_wiki_file_token(meta)
                            if real_token:
                                token = real_token
                                entity_type = "FILE"
                    # 除文件夹外均可通过「导出/下载 + 文本抽取」回退链路导入
                    importable = entity_type != "FOLDER"
                    item = {
                        "doc_token": token,
                        "token": token,
                        "title": title,
                        "url": url,
                        "snippet": snippet,
                        "entity_type": entity_type,
                        "importable": importable,
                    }
                    if not importable:
                        item["unsupported_reason"] = (
                            "文件夹不支持直接导入，请在文件列表中逐个选择文档"
                        )
                    docs.append(item)
                return docs
            # 非标准 ok 格式（如 dws 风格 success）或空 dict → 尝试宽容提取
            if isinstance(result, dict):
                # 尝试常见字段路径
                for key in ("documents", "items", "result", "data"):
                    val = result.get(key)
                    if isinstance(val, list) and val:
                        logger.debug("飞书 doc_search 从 %s 字段提取 %d 条", key, len(val))
                        return val
                    if isinstance(val, dict):
                        sub = val.get("items") or val.get("documents") or val.get("results")
                        if isinstance(sub, list) and sub:
                            return sub
            logger.warning("飞书 doc_search 未从响应中提取到结果: %s",
                        json.dumps(result, ensure_ascii=False)[:300])
            return []
        except IMAdapterPermissionError as e:
            logger.warning("飞书 doc_search 认证失败: %s", e)
            return [{"error": "auth", "message": str(e)}]
        except Exception as e:
            logger.error("飞书 doc_search 异常: %s", e)
            return [{"error": "internal", "message": str(e)}]
        return []

    # ------------------------------------------------------------------
    # 非 docx 文档导入回退链路（导出 / 下载 + 文本抽取）
    # ------------------------------------------------------------------

    def _run_cli_file(self, args: list[str], tmpdir: str,
                      output_base: str = "feishu_doc") -> str:
        """运行会写出本地文件的 CLI 命令（``drive +export`` / ``+download``）。

        lark-cli 要求输出路径为「当前目录内的相对路径」，故用 ``cwd=tmpdir`` 执行，
        再自动定位产出的文件并返回其绝对路径。失败时抛 ``IMAdapter*`` 异常。
        """
        cmd = self._build_command(args, force_no_dry_run=True)
        logger.debug("运行文件型命令: %s (cwd=%s)", " ".join(cmd), tmpdir)
        result = subprocess.run(
            cmd, capture_output=True, text=True,
            encoding="utf-8", env=self._no_browser_env,
            cwd=tmpdir, timeout=self.timeout,
        )
        if result.returncode != 0:
            stderr = (result.stderr or result.stdout or "").strip()
            if result.returncode < 0:
                sig = -result.returncode
                logger.debug("%s 文件型命令被信号 %d 终止（关机阶段）: %s",
                             self.cli_path, sig, stderr)
                raise IMAdapterShutdownError(
                    f"{self.cli_path} terminated by signal {sig}: {stderr}")
            raise self._classify_error(stderr)(
                f"{self.cli_path} exit {result.returncode}: {stderr}")
        files = [
            os.path.join(tmpdir, f) for f in os.listdir(tmpdir)
            if os.path.isfile(os.path.join(tmpdir, f))
        ]
        # 优先取与 output_base 不同名的真实产物（导出分支常同时生成
        # feishu_export 占位 + feishu_export.xlsx 真实文件）；若只剩与
        # output_base 同名的文件（下载分支产出恰好就叫 feishu_dl），则直接取它。
        others = [f for f in files if os.path.basename(f) != output_base]
        produced = others or files
        if not produced:
            raise IMAdapterError(f"{self.cli_path} 未产出文件: {tmpdir}")
        produced.sort(key=lambda p: os.path.getmtime(p), reverse=True)
        return produced[0]

    def _detect_and_rename(self, path: str) -> str:
        """按文件头魔数判断真实格式并在磁盘上重命名（``drive +download`` 拿到的文件无扩展名）。"""
        try:
            with open(path, "rb") as f:
                head = f.read(8)
        except Exception:
            logger.warning("[resilience] silent exception in _detect_and_rename", exc_info=True)
            return path
        ext = None
        if head[:4] == b"%PDF":
            ext = ".pdf"
        elif head[:2] == b"PK":  # docx / xlsx / pptx 均为 zip，按内部目录区分
            try:
                import zipfile
                names = zipfile.ZipFile(path).namelist()
                if any(n.startswith("word/") for n in names):
                    ext = ".docx"
                elif any(n.startswith("xl/") for n in names):
                    ext = ".xlsx"
                elif any(n.startswith("ppt/") for n in names):
                    ext = ".pptx"
                else:
                    ext = ".zip"
            except Exception:
                logger.warning("[resilience] silent exception in _detect_and_rename", exc_info=True)
                ext = ".zip"
        elif head[:8] == b"\xd0\xcf\x11\xe0":  # OLE2（旧版 .doc / .xls）
            ext = ".doc"
        if ext and not path.lower().endswith(ext):
            new_path = path + ext
            try:
                os.rename(path, new_path)
                return new_path
            except OSError as _exc:
                logger.warning(f"_detect_and_rename: swallowed exception: {_exc}")
                return path
        return path  # 兜底：当作纯文本尝试

    def _extract_text_from_file(self, path: str) -> str:
        """从本地文件抽取纯文本，按扩展名选择抽取器。RAG 只关心文本，不要求排版。"""
        lower = path.lower()
        # 文本类：直接读
        if lower.endswith((".md", ".markdown", ".txt", ".text", ".csv", ".tsv",
                           ".json", ".log", ".yaml", ".yml", ".html", ".htm",
                           ".xml", ".py", ".js", ".ts", ".java", ".c", ".cpp",
                           ".go", ".rs", ".sh")):
            with open(path, "r", encoding="utf-8", errors="replace") as f:
                return f.read()
        # Office / PDF 类
        if lower.endswith(".docx"):
            from docx import Document
            doc = Document(path)
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        if lower.endswith(".pdf"):
            import pdfplumber
            with pdfplumber.open(path) as pdf:
                return "\n".join((pg.extract_text() or "") for pg in pdf.pages)
        if lower.endswith((".xlsx", ".xlsm")):
            import openpyxl
            wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
            rows = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) for c in row
                             if c is not None and str(c).strip()]
                    if cells:
                        rows.append("\t".join(cells))
            return "\n".join(rows)
        if lower.endswith(".pptx"):
            from pptx import Presentation
            prs = Presentation(path)
            out = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    # python-pptx stubs don't declare text_frame on BaseShape,
                    # but has_text_frame guarantees it exists at runtime.
                    if shape.has_text_frame and shape.text_frame.text.strip():  # type: ignore[attr-defined]
                        out.append(shape.text_frame.text)  # type: ignore[attr-defined]
            return "\n".join(out)
        # 兜底：当作 UTF-8 文本读取
        try:
            with open(path, "r", encoding="utf-8", errors="replace") as f:
                return f.read()
        except Exception as e:
            raise IMAdapterError(f"无法从文件抽取文本: {path} ({e})") from e

    def _doc_read_fallback(self, node_id: str, entity_type: str | None = None) -> dict:
        """非 docx 类型（file/wiki/sheet/bitable/slides 等）的导入回退链路。

        头尾 RAG 不要求排版，只需纯文本。策略：
        - 可导出类型（doc/docx/wiki/mindnote → markdown；sheet/bitable → xlsx；
          slides → pdf）用 ``drive +export`` 导出为文本友好格式再读取；
        - 通用文件（file）或未知类型用 ``drive +download`` 下载原始文件，
          按魔数识别格式后抽取文本。
        """
        et = (entity_type or "").upper()
        tmpdir = tempfile.mkdtemp(prefix="feishu_doc_")
        try:
            if et in ("FILE", ""):
                # 下载原始文件（file 类型无扩展名，后续按魔数识别）
                path = self._run_cli_file(
                    ["drive", "+download", "--file-token", node_id,
                     "--output", "./feishu_dl", "--overwrite"],
                    tmpdir, output_base="feishu_dl")
                path = self._detect_and_rename(path)
                content = self._extract_text_from_file(path)
            else:
                doc_type = {
                    "DOC": "doc", "DOCX": "docx", "WIKI": "wiki",
                    "SHEET": "sheet", "BITABLE": "bitable",
                    "SLIDE": "slides", "SLIDES": "slides",
                    "MINDNOTE": "mindnote",
                }.get(et)
                if not doc_type:
                    # 未知类型：退化为文件下载
                    path = self._run_cli_file(
                        ["drive", "+download", "--file-token", node_id,
                         "--output", "./feishu_dl", "--overwrite"],
                        tmpdir, output_base="feishu_dl")
                    path = self._detect_and_rename(path)
                    content = self._extract_text_from_file(path)
                else:
                    ext = {
                        "doc": "markdown", "docx": "markdown", "wiki": "markdown",
                        "sheet": "xlsx", "bitable": "xlsx", "slides": "pdf",
                        "mindnote": "markdown",
                    }.get(doc_type, "markdown")
                    path = self._run_cli_file(
                        ["drive", "+export", "--doc-type", doc_type,
                         "--token", node_id, "--file-extension", ext,
                         "--output-dir", ".", "--file-name", "feishu_export",
                         "--overwrite"],
                        tmpdir, output_base="feishu_export")
                    content = self._extract_text_from_file(path)
            return {
                "title": "",  # 标题由调用方用搜索结果补全
                "content": content,
                "doc_type": et or "FILE",
                "url": "",
                "imported_via": "download" if et in ("FILE", "") else "export",
            }
        except IMAdapterShutdownError:
            raise
        except IMAdapterError as e:
            logger.warning("飞书 doc_read 回退失败 (token=%s, type=%s): %s",
                           node_id, et, e)
            return {"error": "extract_failed", "message": str(e), "code": 422}
        except Exception as e:
            logger.error("飞书 doc_read 回退异常 (token=%s, type=%s): %s",
                         node_id, et, e)
            return {"error": "extract_failed", "message": str(e), "code": 500}
        finally:
            shutil.rmtree(tmpdir, ignore_errors=True)

    def doc_read(self, node_id: str, content_format: str = "markdown",
                 entity_type: str | None = None) -> dict:
        """获取飞书文档内容。

        docx 文档走 ``lark-cli docs +fetch`` 拿高质量 markdown；其余类型
        （file/wiki/sheet/bitable/slides 等）走 ``_doc_read_fallback`` 的
        导出 / 下载 + 文本抽取回退链路（头尾 RAG 只关心纯文本，不要求排版）。

        ``node_id`` 即文档 token；``entity_type`` 为搜索结果里的实体类型
        （DOC/DOCX/SHEET/BITABLE/WIKI/FILE/...），用于回退链路分流。

        Returns:
            dict，含 title / content / doc_type / url 等字段；
            异常时返回 ``{"error": "...", "code": ...}`` 供调用方友好展示。
        """
        et = (entity_type or "").upper()
        # 非 docx 类型直接走导出/下载回退（跳过必定失败的 docs +fetch）
        if et not in ("", "DOCX"):
            try:
                return self._doc_read_fallback(node_id, entity_type=et)
            except IMAdapterShutdownError:
                raise
            # 回退内部已把异常转成 error dict，此处直接返回
        args = ["docs", "+fetch", "--doc", node_id,
                "--doc-format", content_format]
        try:
            result = self.run(args, operation="doc_read",
                            force_no_dry_run=True)
            if isinstance(result, dict) and result.get("ok"):
                data = result.get("data", {})
                if isinstance(data, dict):
                    return data
        except IMAdapterPermissionError as e:
            logger.warning("飞书 doc_read 权限错误 (token=%s): %s", node_id, e)
            return {"error": "auth", "message": str(e), "code": 401}
        except IMAdapterResourceNotFoundError as e:
            logger.warning("飞书 doc_read 文档不存在 (token=%s): %s", node_id, e)
            return {"error": "not_found", "message": str(e), "code": 404}
        except IMAdapterRateLimitError as e:
            logger.warning("飞书 doc_read 限频 (token=%s): %s", node_id, e)
            return {"error": "rate_limit", "message": str(e), "code": 429}
        except IMAdapterUnsupportedTypeError:
            # docx 路径失败（如旧 doc 等），再尝试导出/下载回退
            logger.info("飞书 doc_read 类型不支持，转回退链路 (token=%s, type=%s)",
                        node_id, et)
            return self._doc_read_fallback(node_id, entity_type=et or "FILE")
        except Exception as e:
            logger.error("飞书 doc_read 异常 (token=%s): %s", node_id, e)
            return {"error": "internal", "message": str(e), "code": 500}
        return {}

    def doc_list(self, folder_token: str = "", page_size: int = 50,
                 page_token: str = "") -> dict:
        """列出飞书云空间 / 文件夹下的文件清单。

        lark-cli v1.0.72 不存在 ``docs list``；通过 ``drive files list`` 命令
        获取指定文件夹下的文件（含文档）。参数语义保持不变。

        ``folder_token`` 为空时列出根目录文件；支持分页。

        Returns:
            dict，含 items (list[dict]) / has_more / page_token 字段。
            异常时 items 含 error 标记。
        """
        args = ["drive", "files", "list", "--page-size", str(page_size)]
        if folder_token:
            args += ["--folder-token", folder_token]
        if page_token:
            args += ["--page-token", page_token]
        try:
            result = self.run(args, operation="doc_list",
                            force_no_dry_run=True)
            if isinstance(result, dict) and result.get("ok"):
                data = result.get("data", {})
                items = data.get("files") or data.get("items") or data.get("data") or []
                return {
                    "items": items,
                    "has_more": data.get("has_more", False),
                    "page_token": data.get("page_token", ""),
                }
        except IMAdapterPermissionError as e:
            logger.warning("飞书 doc_list 认证失败: %s", e)
            return {"items": [{"error": "auth", "message": str(e)}],
                    "has_more": False, "page_token": ""}
        except Exception as e:
            logger.error("飞书 doc_list 异常: %s", e)
            return {"items": [{"error": "internal", "message": str(e)}],
                    "has_more": False, "page_token": ""}
        return {"items": [], "has_more": False, "page_token": ""}

    # ------------------------------------------------------------------

