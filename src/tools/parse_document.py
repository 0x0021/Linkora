"""文档解析工具：支持 PDF、PPT、Word、图片（OCR）"""
from __future__ import annotations

import logging
import re
from pathlib import Path
from typing import Optional

from src.config import AppConfig

logger = logging.getLogger(__name__)

# ============================================================================
# OCR 后处理：去噪 + 格式整理
# ============================================================================
# RapidOCR 逐行输出，聊天截图里混有大量 UI 噪点（按钮文字、时间戳、页码）、
# 重复行、以及被拆散的「标签 / 值」对（如「总金额（元）」「103.791.68」两行）。
# 这些脏内容直接喂给 LLM 会严重干扰理解。post_process_ocr_text 在落库 / 投喂
# LLM / 历史上下文读取三处统一使用，保证全链路干净。

# 通用 UI 外壳词（任意 App 截图都可能出现，整行完全等于这些词时视为噪点丢弃）。
# 注意：仅整行精确匹配才丢弃，含这些词的正常句子（如「查看详情」）会保留。
_UI_CHROME_TOKENS: frozenset[str] = frozenset({
    # 导航 / 操作
    "返回", "关闭", "取消", "确定", "完成", "提交", "下一步", "上一步", "保存",
    "编辑", "删除", "复制", "转发", "收藏", "分享", "看分布", "详情", "查看", "立即查看",
    "更多", "收起", "展开", "刷新", "搜索", "添加", "下载", "上传", "设置",
    "帮助", "客服", "举报", "投诉", "登录", "注册", "退出", "首页", "尾页",
    "上一页", "下一页", "全部", "暂无数据", "加载更多", "广告", "推荐", "热门",
    # 社交 / 互动
    "点赞", "评论", "关注", "主页", "消息", "我的", "表情", "语音", "图片",
    "通讯录", "工作台", "日历", "文档", "扫一扫", "发送",
    # 钉钉 / 通用品牌与状态
    "钉钉", "AI", "loading", "Loading", "加载中", "正在加载",
    # 装饰符号 / 省略
    "···", "...", "..", "+", "＋", ">", "›", "»", "·", "•",
})

# 行尾常见的 OCR 误识标点（识别错误的引号 / 书名号 / 箭头），清理掉更干净
_TRAILING_JUNK_RE = re.compile(r"[》”’」』「『>»·•\-=|]+$")
_LEADING_JUNK_RE = re.compile(r"^[》”’」』「『>»·•\s]+")
# 零宽 / 不可见字符
_ZERO_WIDTH_RE = re.compile(r"[\u200b\u200c\u200d\ufeff\u00a0]")
_WS_RE = re.compile(r"\s+")                       # 任意空白折叠为单空格
# 时钟碎屑（状态栏时间 / 消息时间戳，如 21:05 / 9:30 / 21:051）
_CLOCK_RE = re.compile(r"^\d{1,2}[:：]\d{1,4}$")
_PURE_SYM_RE = re.compile(r"^[\W_]+$")            # 纯符号 / 标点行
# 列表 / 条目标记（数字序号、项目符号），作为段落边界
_LIST_MARK_RE = re.compile(r"^\s*(\d{1,3}[.、)]\s|[（(]\d+[)）]|[•▪◦\-*]\s)")

_SENT_END = set("。！？.!?…")                       # 句末标点 → 新段落
_LABEL_END = ("：", ":", "（", "(", "）", ")")      # 标签行结尾 → 新段落
_DIGIT_START = set("0123456789¥$￥+-（(")           # 值行开头
_CJK_RE = re.compile(r"[\u4e00-\u9fff]")


def _is_noise_line(line: str) -> bool:
    """判断单行是否为 OCR 噪点，应丢弃（保守、通用）。"""
    s = line.strip()
    if not s:
        return True
    # 纯符号 / 标点行
    if _PURE_SYM_RE.match(s):
        return True
    # 整行等于通用 UI 外壳词
    if s in _UI_CHROME_TOKENS:
        return True
    # 独立时钟碎屑（状态栏 / 时间戳）
    if _CLOCK_RE.match(s):
        return True
    return False


def _looks_like_label(line: str) -> bool:
    """是否为「标签行」：短描述、自身不含数字、且带标签特征。

    采用结构判断而非业务词表，对表单 / 票据 / 配置页等各类截图通用。
    判定偏严：仅显式以冒号 / 左括号结尾，或极短（≤ 4 字）名词短语视为标签，
    避免把普通短句误判为字段名。
    """
    s = line.strip()
    if not s or any(ch.isdigit() for ch in s):
        return False
    if s in _UI_CHROME_TOKENS:
        return False
    # 显式以冒号 / 左括号结尾（如「总金额（元）」「姓名：」）→ 标签
    if s.endswith(_LABEL_END):
        return True
    # 极短名词短语（≤ 4 字、无句末标点）→ 多半是字段名
    if len(s) <= 4 and s[-1] not in _SENT_END:
        return True
    return False


def _is_merge_value(nxt: str) -> bool:
    """是否为可并入上一「标签行」的「值行」，判定偏严。"""
    if not nxt:
        return False
    # 列表项不是值
    if _LIST_MARK_RE.match(nxt):
        return False
    # 自带标签结构（含冒号 / 括号）的说明这行自己就是「标签 + 内容」
    if "：" in nxt or ":" in nxt or "（" in nxt or "(" in nxt:
        return False
    # 下一行自身像标签（如另一个字段名）→ 不并入
    if _looks_like_label(nxt):
        return False
    return _looks_like_value(nxt)


def _looks_like_value(line: str) -> bool:
    """是否为「值行」：以数字 / 货币 / 正负号 / 左括号开头，或含数字。"""
    s = line.strip()
    if not s:
        return False
    if s[0] in _DIGIT_START:
        return True
    return any(ch.isdigit() for ch in s)


def _clean_line(s: str) -> str:
    """行级基础清洗：去零宽字符、去首尾误识标点、折叠空白。"""
    s = _ZERO_WIDTH_RE.sub("", s)
    s = _LEADING_JUNK_RE.sub("", s)
    s = _TRAILING_JUNK_RE.sub("", s)
    s = _WS_RE.sub(" ", s).strip()
    return s


def _should_join(prev: str, nxt: str) -> bool:
    """判断 prev 与 nxt 是否应合并为同一段落（软换行修复）。

    返回 False 的典型边界：标签行结尾、下一条是列表 / 数字开头、下一行是短
    标签（说话人名 / 小标题）、上一行以句末标点结束。
    """
    if not prev or not nxt:
        return False
    if prev.endswith(_LABEL_END):
        return False
    if _LIST_MARK_RE.match(nxt) or nxt[0] in _DIGIT_START:
        return False
    # 显式以冒号结尾的下一行（新字段）→ 新段落
    if nxt.endswith(("：", ":")):
        return False
    # 下一行自带标签结构（含冒号 / 括号）→ 新字段，保留边界
    if "：" in nxt or ":" in nxt or "（" in nxt or "(" in nxt:
        return False
    # 下一行以右括号结尾（如「总金额（元）」）→ 新字段，保留边界
    if nxt.endswith(("）", ")")):
        return False
    # 短 token（≤ 5 字）后接中文整句：多为说话人名 / 小标题，保留边界
    if len(prev) <= 5 and _CJK_RE.match(nxt[0]):
        return False
    # 短下一行（≤ 4 字、无句末标点）：多为说话人名 / 小标题，其前保留边界
    if len(nxt) <= 4 and nxt[-1] not in _SENT_END:
        return False
    if prev[-1] in _SENT_END:
        return False
    return True


def post_process_ocr_text(raw_text: str, max_chars: int = 1500) -> str:
    """对 OCR 原始文本做通用后处理，产出干净、连贯、紧凑的文本。

    处理流程（通用、不假设具体业务）：
      1. 行级清洗：去零宽字符、去噪点行（纯符号 / UI 外壳 / 时钟碎屑）
      2. 去连续重复行
      3. 结构重组：紧邻的「标签行 + 值行」合并为「标签：值」
      4. 软换行修复：按标点 / 标签 / 列表结构把被切散的正文重新拼成段落
      5. 超长截断（默认保留前 1500 字，避免撑爆上下文）

    Args:
        raw_text: RapidOCR 逐行拼接的原始文本
        max_chars: 截断阈值，超过则保留前 max_chars 字并追加提示

    Returns:
        处理后文本；若处理后为空（原始即噪声），回退返回原始 strip 结果，
        避免信息彻底丢失。
    """
    if not raw_text or not raw_text.strip():
        return ""

    # 1. 行级清洗 + 去噪
    lines: list[str] = []
    for ln in raw_text.split("\n"):
        c = _clean_line(ln)
        if c and not _is_noise_line(c):
            lines.append(c)

    # 2. 去连续重复行
    dedup: list[str] = []
    for ln in lines:
        if dedup and dedup[-1] == ln:
            continue
        dedup.append(ln)

    if not dedup:
        # 全是噪声：回退到原始，至少保留信息
        return raw_text.strip()

    # 3. 标签 + 值 合并（结构判断，通用）
    merged: list[str] = []
    i = 0
    n = len(dedup)
    while i < n:
        cur = dedup[i]
        nxt = dedup[i + 1] if i + 1 < n else None
        if nxt is not None and _looks_like_label(cur) and _is_merge_value(nxt):
            merged.append(f"{cur.rstrip('：:')}：{nxt}")
            i += 2
            continue
        merged.append(cur)
        i += 1

    # 4. 软换行修复：把被切散的正文拼回段落
    paras: list[str] = []
    buf = merged[0]
    for k in range(1, len(merged)):
        nxt = merged[k]
        if _should_join(buf, nxt):
            sep = " " if (buf[-1].isascii() and nxt[0].isascii()) else ""
            buf = buf + sep + nxt
        else:
            paras.append(buf)
            buf = nxt
    paras.append(buf)

    text = "\n".join(p for p in paras if p.strip())

    # 5. 超长截断
    if len(text) > max_chars:
        text = text[:max_chars].rstrip() + f"\n…[OCR 内容过长已截断，原文共 {len(text)} 字]"

    return text

# RapidOCR 针对「钉钉聊天截图」场景调优的参数（覆盖默认 PP-OCRv3 配置）。
# 默认配置对截图偏保守（min_height=30 / limit_side_len=736 会丢小字与大图细节），
# 以下参数在保持速度的前提下显著提升小字号、长行、高清截图的识别召回率。
OCR_TUNED_PARAMS: dict = {
    "text_score": 0.4,           # 降置信阈值，召回低对比/小字
    "min_height": 15,            # 默 30 太高，截图里大量小字号
    "limit_side_len": 1080,      # 默 736 会降采样高清图，改 1080 保留细节
    "limit_type": "min",
    "unclip_ratio": 1.8,         # 文本框更贴合，减少截断
    "width_height_ratio": 10,    # 允许更长行
    "box_thresh": 0.45,          # 略降，配合 text_score 提召回
    "max_candidates": 2000,      # 更多候选框
    "use_dilation": True,
}


class DocumentParser:
    """统一文档解析接口，支持多种格式。"""

    def __init__(self, config: AppConfig):
        self.config = config
        self._pdf_available = False
        self._ppt_available = False
        self._docx_available = False
        self._ocr_available = False
        self._check_dependencies()

    def _check_dependencies(self) -> None:
        """检查各解析器依赖是否可用。"""
        # PDF
        try:
            import pdfplumber  # noqa: F401  (import 即可用性检查)
            self._pdf_available = True
            logger.info("[解析器] PDF 解析可用（pdfplumber）")
        except ImportError:
            logger.warning("[解析器] PDF 解析不可用，请安装: pip install pdfplumber")

        # PPT
        try:
            import pptx  # noqa: F401  (import 即可用性检查)
            self._ppt_available = True
            logger.info("[解析器] PPT 解析可用（python-pptx）")
        except ImportError:
            logger.warning("[解析器] PPT 解析不可用，请安装: pip install python-pptx")

        # Word
        try:
            import docx  # noqa: F401  (import 即可用性检查)
            self._docx_available = True
            logger.info("[解析器] Word 解析可用（python-docx）")
        except ImportError:
            logger.warning("[解析器] Word 解析不可用，请安装: pip install python-docx")

        # OCR（使用 RapidOCR，PP-OCRv3 模型，中文识别远超 Tesseract）
        try:
            from rapidocr_onnxruntime import RapidOCR
            # 预创建引擎实例，避免首次调用时下载模型阻塞；
            # 注入针对截图调优的参数（见模块级 OCR_TUNED_PARAMS）
            self._ocr_engine = RapidOCR(**OCR_TUNED_PARAMS)
            self._ocr_available = True
            logger.info("[解析器] OCR 识别可用（RapidOCR + PP-OCRv3，已应用截图调优参数）")
        except ImportError:
            logger.warning("[解析器] OCR 不可用，请安装: pip install rapidocr-onnxruntime")
            self._ocr_engine = None
            self._ocr_available = False

    def parse(self, file_path: str, file_type: Optional[str] = None) -> str:
        """解析文档，返回纯文本。

        Args:
            file_path: 文件路径
            file_type: 文件类型（可选，自动从扩展名推断）

        Returns:
            提取的纯文本（失败返回空字符串）
        """
        path = Path(file_path)
        if not path.exists():
            logger.error("[解析器] 文件不存在: %s", file_path)
            return ""

        # 推断文件类型
        if file_type is None:
            suffix = path.suffix.lower()
            if suffix == ".pdf":
                file_type = "pdf"
            elif suffix in (".ppt", ".pptx"):
                file_type = "ppt"
            elif suffix in (".doc", ".docx"):
                file_type = "docx"
            elif suffix in (".png", ".jpg", ".jpeg", ".bmp", ".tiff"):
                file_type = "image"
            elif suffix in (".html", ".htm"):
                file_type = "html"
            elif suffix in (".md", ".markdown", ".mdown"):
                file_type = "markdown"
            elif suffix in (".txt", ".text"):
                file_type = "text"
            else:
                logger.error("[解析器] 不支持的文件类型: %s", suffix)
                return ""

        # 调用对应解析器
        if file_type == "pdf":
            return self._parse_pdf(file_path)
        elif file_type == "ppt":
            return self._parse_ppt(file_path)
        elif file_type == "docx":
            return self._parse_docx(file_path)
        elif file_type == "image":
            return self._parse_image(file_path)
        elif file_type == "html":
            return self._parse_html(file_path)
        elif file_type in ("markdown", "text"):
            return self._parse_markdown(file_path)
        else:
            logger.error("[解析器] 未知的文件类型: %s", file_type)
            return ""

    def _parse_pdf(self, file_path: str) -> str:
        """解析 PDF 文件。

        策略：
        1. 使用 pdfplumber 提取文本
        2. 过滤页眉页脚（通过检测每页顶部/底部的重复文本）
        3. 合并页面内容并清理格式
        """
        if not self._pdf_available:
            logger.error("[解析器] PDF 解析依赖未安装")
            return ""

        try:
            import pdfplumber
            text_parts = []
            page_texts = []

            with pdfplumber.open(file_path) as pdf:
                for _, page in enumerate(pdf.pages):
                    page_text = page.extract_text()
                    if page_text:
                        page_texts.append(page_text)

            if not page_texts:
                logger.warning("[解析器] PDF 无文本内容，可能是扫描版，尝试 OCR...")
                return self._parse_pdf_ocr(file_path)

            # 过滤页眉页脚：统计所有页面的前3行和后3行，找出重复出现的
            header_candidates = {}
            footer_candidates = {}

            for page_text in page_texts:
                lines = [line.strip() for line in page_text.split('\n') if line.strip()]
                if len(lines) >= 6:
                    # 前3行作为页眉候选
                    for j in range(min(3, len(lines))):
                        line = lines[j]
                        header_candidates[line] = header_candidates.get(line, 0) + 1
                    # 后3行作为页脚候选
                    for j in range(max(0, len(lines)-3), len(lines)):
                        line = lines[j]
                        footer_candidates[line] = footer_candidates.get(line, 0) + 1

            # 出现在超过50%页面的行视为页眉页脚
            threshold = max(2, len(page_texts) * 0.5)
            headers_to_remove = {k for k, v in header_candidates.items() if v >= threshold}
            footers_to_remove = {k for k, v in footer_candidates.items() if v >= threshold}

            # 提取每页正文，去除页眉页脚
            for page_text in page_texts:
                lines = [line.strip() for line in page_text.split('\n') if line.strip()]
                # 过滤页眉页脚
                filtered_lines = [
                    line for line in lines
                    if line not in headers_to_remove and line not in footers_to_remove
                ]
                if filtered_lines:
                    text_parts.append('\n'.join(filtered_lines))

            result = '\n\n'.join(text_parts)

            # 如果提取的文本太少，可能是纯图片PDF，尝试OCR
            if len(result.strip()) < 100 and self._ocr_available:
                logger.info("[解析器] PDF 文本过少，尝试 OCR 识别...")
                ocr_result = self._parse_pdf_ocr(file_path)
                if ocr_result:
                    return ocr_result

            logger.info("[解析器] PDF 解析完成: %s（%d 页，提取 %d 字符）",
                       file_path, len(page_texts), len(result))
            return result
        except Exception as e:
            logger.error("[解析器] PDF 解析失败: %s", e, exc_info=True)
            return ""

    def _parse_pdf_ocr(self, file_path: str) -> str:
        """使用 OCR 解析 PDF（适用于扫描版 PDF）。"""
        if not self._ocr_available or self._ocr_engine is None:
            logger.warning("[解析器] OCR 不可用，无法解析扫描版 PDF")
            return ""

        try:
            import fitz  # PyMuPDF
            import tempfile

            doc = fitz.open(file_path)
            try:
                text_parts = []

                for page_num in range(len(doc)):
                    page = doc.load_page(page_num)
                    pix = page.get_pixmap(dpi=300)
                    img_data = pix.tobytes("png")
                    # 写入临时文件供 RapidOCR 读取
                    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
                        tmp.write(img_data)
                        tmp_path = tmp.name
                    try:
                        result, elapse = self._ocr_engine(tmp_path)
                        if result:
                            lines = [item[1] for item in result if item[1] and item[1].strip()]
                            page_text = "\n".join(lines)
                            if page_text.strip():
                                text_parts.append(f"[第 {page_num + 1} 页]\n{page_text.strip()}")
                    finally:
                        import os
                        try:
                            os.unlink(tmp_path)
                        except Exception as e:
                            logger.debug("unlink %s failed: %s", tmp_path, e)

            finally:
                doc.close()
            result = '\n\n'.join(text_parts)

            if result.strip():
                logger.info("[解析器] PDF OCR 解析完成: %s（提取 %d 字符）",
                           file_path, len(result))
                return result
            else:
                logger.warning("[解析器] PDF OCR 未识别到有效内容")
                return ""
        except ImportError:
            logger.warning("[解析器] PDF OCR 依赖缺失（缺少 PyMuPDF），无法解析扫描版 PDF")
            return ""
        except Exception as e:
            logger.error("[解析器] PDF OCR 解析失败: %s", e, exc_info=True)
            return ""

    def _parse_ppt(self, file_path: str) -> str:
        """解析 PPT 文件（.pptx）。"""
        if not self._ppt_available:
            logger.error("[解析器] PPT 解析依赖未安装")
            return ""

        try:
            from pptx import Presentation
            prs = Presentation(file_path)

            text_parts = []
            for i, slide in enumerate(prs.slides):
                slide_texts = []
                for shape in slide.shapes:
                    text = getattr(shape, "text", None)
                    if text:
                        slide_texts.append(text)

                if slide_texts:
                    text_parts.append(f"[第 {i+1} 页]\n" + "\n".join(slide_texts))

            result = "\n\n".join(text_parts)
            logger.info("[解析器] PPT 解析完成: %s（%d 页，提取 %d 字符）",
                        file_path, len(prs.slides), len(result))
            return result
        except Exception as e:
            logger.error("[解析器] PPT 解析失败: %s", e)
            return ""

    def _parse_docx(self, file_path: str) -> str:
        """解析 Word 文件（.docx）。"""
        if not self._docx_available:
            logger.error("[解析器] Word 解析依赖未安装")
            return ""

        try:
            from docx import Document
            doc = Document(file_path)

            text_parts = []
            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)

            # 也提取表格内容
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            text_parts.append(cell.text)

            result = "\n".join(text_parts)
            logger.info("[解析器] Word 解析完成: %s（提取 %d 字符）", file_path, len(result))
            return result
        except Exception as e:
            logger.error("[解析器] Word 解析失败: %s", e)
            return ""

    def _parse_html(self, file_path: str) -> str:
        """解析 HTML 文件，提取可见正文内容。

        策略：
        1. 首先检测页面类型（导航页/文章页/列表页等）
        2. 根据页面类型选择合适的提取策略
        3. 清理脚本、样式等无关元素
        4. 过滤重复和无意义的内容
        """
        try:
            from bs4 import BeautifulSoup

            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                html_content = f.read()

            soup = BeautifulSoup(html_content, "html.parser")

            # 检测页面类型
            page_type = self._detect_page_type(soup)
            logger.info("[解析器] 检测到页面类型: %s", page_type)

            if page_type == "navigation":
                return self._extract_navigation_content(soup)
            elif page_type == "article":
                return self._extract_article_content(soup)
            else:
                return self._extract_general_content(soup)
        except Exception as e:
            logger.error("[解析器] HTML 解析失败: %s", e, exc_info=True)
            return ""

    def _detect_page_type(self, soup) -> str:
        """检测页面类型：navigation/article/list/general"""

        # 获取标题和meta信息
        title = soup.find('title')
        title_text = title.get_text().lower() if title else ''

        meta_desc = soup.find('meta', attrs={'name': 'description'})
        meta_desc.get('content', '').lower() if meta_desc else ''

        # 统计链接数量
        links = soup.find_all('a', href=True)
        link_count = len(links)

        # 统计文本块数量
        text_blocks = [p.get_text(strip=True) for p in soup.find_all(['p', 'li', 'h1', 'h2', 'h3'])]
        text_block_count = len([t for t in text_blocks if len(t) > 10])

        # 判断是否为导航页：链接较多且包含业务链接（非仅导航栏）
        # 条件：链接数>=5 或 标题包含"导航"关键词
        is_nav_title = any(kw in title_text for kw in ['导航', 'directory', 'index', 'portal'])
        has_business_links = any(
            link.get('href', '').startswith(('http://', 'https://')) and
            not self._is_navigation_link(link.get_text(strip=True), link.get('href', ''))
            for link in links[:20]  # 只检查前20个链接
        )

        if (link_count >= 5 or is_nav_title) and has_business_links:
            return 'navigation'

        # 判断是否为文章页：有article/main标签或大量段落
        if soup.find(['article', 'main']) or text_block_count > 10:
            return 'article'

        # 判断是否为列表页
        if soup.find(['ul', 'ol']) and link_count > 5:
            return 'list'

        return 'general'

    def _extract_navigation_content(self, soup) -> str:
        """提取导航页内容：保留网站名称、描述和链接信息

        导航页的特点是所有链接都是主要内容，不应过滤
        """

        # 移除脚本和样式
        for tag in soup.find_all(['script', 'style', 'noscript']):
            tag.decompose()

        content_parts = []

        # 1. 提取页面标题
        title = soup.find('title')
        if title:
            title_text = title.get_text(strip=True)
            if title_text:
                content_parts.append(f"# {title_text}")

        # 2. 提取页面主标题（h1, h2）
        for heading in soup.find_all(['h1', 'h2']):
            text = heading.get_text(strip=True)
            if text and len(text) > 2:
                # 过滤明显的UI元素文本
                if not self._is_ui_element_text(text):
                    prefix = "## " if heading.name == 'h2' else "# "
                    content_parts.append(f"{prefix}{text}")

        # 3. 提取所有有意义的链接（导航页中所有链接都是主要内容）
        links_data = []
        processed_links = set()  # 避免重复
        link_descriptions = set()  # 记录已提取的描述，避免重复

        for link in soup.find_all('a', href=True):
            href = link.get('href', '').strip()
            text = link.get_text(strip=True)

            # 跳过空链接或javascript链接
            if not text or not href or href.startswith('javascript:') or href.startswith('#'):
                continue

            # 去重
            link_key = f"{text}|{href}"
            if link_key in processed_links:
                continue
            processed_links.add(link_key)

            # 获取链接的描述
            description = self._get_link_description(link)

            # 格式化链接信息
            link_info = f"- **{text}**: [{href}]({href})"
            if description:
                link_info += f" - {description}"
                link_descriptions.add(description)  # 记录已使用的描述

            links_data.append(link_info)

        if links_data:
            content_parts.append("\n## 系统/页面列表")
            content_parts.extend(links_data)

        # 4. 提取有意义的段落和描述（排除已在链接中出现的描述）
        for p in soup.find_all('p'):
            text = p.get_text(strip=True)
            # 保留较长的描述性文本
            if text and len(text) > 15:
                # 过滤明显的UI元素文本
                if self._is_ui_element_text(text):
                    continue
                # 过滤已经在链接描述中出现的文本，避免重复
                if text in link_descriptions:
                    continue
                # 过滤版权信息和备案号
                if any(kw in text.lower() for kw in ['©', 'copyright', 'icp备', '版权所有']):
                    continue
                content_parts.append(text)

        result = '\n\n'.join(content_parts)
        logger.info("[解析器] 导航页解析完成: 提取 %d 字符，%d 个链接",
                   len(result), len(links_data))
        return result

    def _extract_article_content(self, soup) -> str:
        """提取文章内容：传统文章的解析逻辑"""

        # 移除不需要的标签
        tags_to_remove = ['script', 'style', 'noscript', 'iframe', 'head',
                         'nav', 'header', 'footer', 'aside',
                         'menu', 'toolbar', 'ad', 'advertisement',
                         'button', 'form', 'input', 'select', 'textarea']
        for tag_name in tags_to_remove:
            for tag in soup.find_all(tag_name):
                tag.decompose()

        # 提取主要内容区域
        main_content = None

        # 优先级：main > article > content类 > body > soup
        for selector in [
            ('main', None),
            ('article', None),
            ('div', lambda x: x and any(p in x.lower() for p in ['content', 'main', 'article', 'body', 'post', 'document'])),
            ('section', lambda x: x and any(p in x.lower() for p in ['content', 'main', 'article', 'body'])),
            ('body', None)
        ]:
            tag_name, attr_filter = selector
            if attr_filter:
                target = soup.find(tag_name, id=attr_filter)
                if not target:
                    target = soup.find(tag_name, class_=attr_filter)
            else:
                target = soup.find(tag_name)

            if target:
                main_content = target
                break

        if not main_content:
            main_content = soup

        # 提取文本并清理
        text = main_content.get_text(separator='\n', strip=True)
        lines = []
        seen_lines = set()

        for line in text.splitlines():
            line = line.strip()
            if not line or len(line) < 3 or line in seen_lines:
                continue

            # 过滤导航文本
            if self._is_navigation_text(line):
                continue

            seen_lines.add(line)
            lines.append(line)

        result = '\n'.join(lines)
        logger.info("[解析器] 文章页解析完成: 提取 %d 字符，%d 行",
                   len(result), len(lines))
        return result

    def _extract_general_content(self, soup) -> str:
        """提取通用页面内容"""

        # 移除脚本和样式
        for tag in soup.find_all(['script', 'style', 'noscript', 'iframe']):
            tag.decompose()

        # 提取所有有意义的文本
        text_parts = []

        # 标题
        for heading in soup.find_all(['h1', 'h2', 'h3', 'h4']):
            text = heading.get_text(strip=True)
            if text and len(text) > 2:
                text_parts.append(text)

        # 段落
        for p in soup.find_all('p'):
            text = p.get_text(strip=True)
            if text and len(text) > 20:
                text_parts.append(text)

        # 列表项
        for li in soup.find_all('li'):
            text = li.get_text(strip=True)
            if text and len(text) > 10:
                text_parts.append(text)

        result = '\n\n'.join(text_parts)
        logger.info("[解析器] 通用页解析完成: 提取 %d 字符", len(result))
        return result

    def _is_navigation_text(self, text: str) -> bool:
        """判断文本是否为导航性文本"""
        import re

        nav_patterns = [
            r'^首页$', r'^产品$', r'^服务$', r'^文档$', r'^登录$', r'^注册$',
            r'^退出$', r'^搜索$', r'^联系我们$', r'^关于我们$',
            r'^home$', r'^login$', r'^register$', r'^search$', r'^menu$',
            r'^logout$', r'^contact$', r'^about$', r'^help$', r'^faq$',
            r'跳转到主内容', r'skip to', r'breadcrumb', r'面包屑',
            r'©.*版权所有', r'copyright', r'备案号', r'icp备',
            r'^帮助中心$', r'^常见问题$',
            r'未找到匹配的系统',  # 空状态提示
        ]

        for pattern in nav_patterns:
            if re.search(pattern, text, re.IGNORECASE):
                return True
        return False

    def _is_ui_element_text(self, text: str) -> bool:
        """判断文本是否为UI元素文本（按钮、操作等）"""

        # UI元素关键词列表（包含匹配）
        ui_keywords = [
            '关闭', '取消', '确定', '提交', '保存',
            '删除', '编辑', '新增', '创建',
            'close', 'cancel', 'ok', 'submit', 'save',
            'delete', 'edit', 'add', 'create',
            '跳转到主内容', 'skip to', 'breadcrumb', '面包屑',
            '未找到匹配', 'not found', 'no results',
            '×', '✕', '···', '...',  # 特殊符号
        ]

        # 检查是否包含UI关键词
        for keyword in ui_keywords:
            if keyword and keyword.lower() in text.lower():
                return True

        # 检查纯符号或短文本（长度<=3且不含字母数字）
        if len(text.strip()) <= 3 and not any(c.isalnum() for c in text):
            return True

        return False

    def _is_navigation_link(self, text: str, href: str) -> bool:
        """判断链接是否为导航性链接"""
        import re

        # 常见的导航链接模式
        nav_href_patterns = [
            r'^/$', r'^/#', r'^/login', r'^/register', r'^/search',
            r'^javascript:', r'^mailto:', r'tel:'
        ]

        nav_text_patterns = [
            r'^首页$', r'^主页$', r'^登录$', r'^注册$', r'^退出$',
            r'^搜索$', r'^设置$', r'^个人中心$', r'^我的',
            r'^home$', r'^login$', r'^signup$', r'^signout$',
            r'^settings$', r'^profile$'
        ]

        for pattern in nav_href_patterns:
            if re.search(pattern, href, re.IGNORECASE):
                return True

        for pattern in nav_text_patterns:
            if re.search(pattern, text, re.IGNORECASE):
                return True

        return False

    def _get_link_description(self, link) -> str:
        """获取链接的描述文本"""
        # 尝试从父元素的下一个兄弟元素获取描述
        parent = link.parent
        if parent:
            next_sibling = parent.find_next_sibling()
            if next_sibling and next_sibling.name in ['p', 'span', 'div']:
                desc = next_sibling.get_text(strip=True)
                if desc and len(desc) < 200:
                    return desc

        # 尝试从link的title属性获取
        title = link.get('title', '')
        if title and len(title) < 200:
            return title

        # 尝试从aria-label获取
        aria_label = link.get('aria-label', '')
        if aria_label and len(aria_label) < 200:
            return aria_label

        return ''

    def _parse_image(self, file_path: str) -> str:
        """OCR 识别图片中的文字（使用 RapidOCR PP-OCRv4 模型）。"""
        if not self._ocr_available or self._ocr_engine is None:
            logger.error("[解析器] OCR 依赖未安装")
            return ""

        try:
            result, elapse = self._ocr_engine(file_path)
            if not result:
                logger.warning("[解析器] OCR 未识别到文字: %s", file_path)
                return ""

            # 按阅读顺序拼接各行文本
            lines = [item[1] for item in result if item[1] and item[1].strip()]
            text = "\n".join(lines)
            if text.strip():
                logger.info("[解析器] OCR 识别完成: %s（提取 %d 字符，耗时 %.2fs）",
                            file_path, len(text), elapse[0] if elapse else 0)
                return text.strip()
            return ""
        except Exception as e:
            logger.error("[解析器] OCR 识别失败: %s", e)
            return ""

    def ocr_image(self, file_path: str) -> str:
        """公开 OCR 接口：识别图片中的文字（供聊天图片识别调用）。"""
        return self._parse_image(file_path)


    def _parse_markdown(self, file_path: str) -> str:
        """解析 Markdown / 纯文本文件。"""
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                content = f.read()

            # 如果是 Markdown，可选择性保留结构（标题、列表等）
            # 这里直接返回原始文本，保留换行符
            result = content.strip()

            if result:
                logger.info("[解析器] Markdown/文本解析完成: %s（提取 %d 字符）",
                           file_path, len(result))
                return result
            else:
                logger.warning("[解析器] Markdown/文本为空: %s", file_path)
                return ""
        except Exception as e:
            logger.error("[解析器] Markdown/文本解析失败: %s", e, exc_info=True)
            return ""


